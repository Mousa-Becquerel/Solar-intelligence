import json
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re
from datetime import datetime
import shutil

def load_json_data(json_file_path):
    """Load and parse the JSON conversation data"""
    with open(json_file_path, 'r') as file:
        data = json.load(file)
    
    plot_items = []
    for item in data['items']:
        if item['type'] == 'plot':
            plot_items.append(item)
    
    if not plot_items:
        raise ValueError("No plot data found in JSON file")
    
    return plot_items, data

def prepare_chart_data_generic(plot_data):
    """Convert plot data to format suitable for PowerPoint chart"""
    plot_type = plot_data.get('plot_type', 'unknown')
    data = plot_data.get('data', [])
    
    if plot_type == 'line':
        # Get all unique years and series names
        years = sorted(list(set([item['date'].split('-')[0] for item in data if 'date' in item])))
        series_names = list(set([item.get('series', 'Data') for item in data if 'series' in item]))
        
        # Organize data by series
        series_data = {}
        for series_name in series_names:
            series_data[series_name] = {}
        
        for item in data:
            if 'date' in item and 'series' in item:
                year = item['date'].split('-')[0]
                series = item['series']
                value = item.get('value', 0)
                series_data[series][year] = value
        
        # Convert to ordered lists for each series
        series_values = {}
        for series_name in series_names:
            values = []
            for year in years:
                values.append(series_data[series_name].get(year, 0))
            series_values[series_name] = values
        
        # Calculate totals (using first series for compatibility)
        first_series = series_names[0] if series_names else 'Data'
        totals = {year: series_data.get(first_series, {}).get(year, 0) for year in years}
        
        return {
            'type': 'line',
            'categories': years,
            'series': series_values,
            'series_names': series_names,
            'years': years,
            'totals': totals
        }
    
    elif plot_type == 'stacked':
        years = sorted(list(set([item.get('category') for item in data if 'category' in item])))
        series_names = list(set([item.get('series') for item in data if 'series' in item]))

        chart_data = {}
        total_by_year = {}

        for item in data:
            year = item.get('category')
            series = item.get('series')
            value = item.get('value', 0)

            if year is not None and series is not None:
                if year not in chart_data:
                    chart_data[year] = {}
                    total_by_year[year] = 0

                chart_data[year][series] = value
                total_by_year[year] += value

        return {
            'type': 'stacked',
            'categories': [str(year) for year in years],
            'series': chart_data,
            'series_names': series_names,
            'years': years,
            'totals': total_by_year
        }

    elif plot_type == 'pie':
        # Extract categories and values from pie chart data
        categories = []
        values = []

        for item in data:
            category = item.get('category') or item.get('label') or item.get('name')
            value = item.get('value', 0)

            if category is not None and value is not None:
                categories.append(str(category))
                values.append(value)

        total = sum(values)

        return {
            'type': 'pie',
            'categories': categories,
            'values': values,
            'total': total
        }

    else:
        return {
            'type': 'generic',
            'categories': [],
            'series': {},
            'years': [],
            'totals': {}
        }

def generate_generic_content(plot_data, chart_info):
    """Generate generic content that can be easily modified later"""
    title = plot_data.get('title', 'Data Analysis Chart')
    unit = plot_data.get('unit', 'Units')
    
    if chart_info['type'] == 'line':
        years = chart_info['years']
        series_names = chart_info.get('series_names', [])
        
        if years and len(years) >= 2:
            # Use first series for main calculations
            first_series = series_names[0] if series_names else 'Data'
            if first_series in chart_info['series']:
                values = chart_info['series'][first_series]
                start_val = values[0] if values else 0
                end_val = values[-1] if values else 0
                growth = ((end_val - start_val) / start_val * 100) if start_val > 0 else 0
                
                # Create series summary
                series_summary = f"• Scenarios analyzed: {', '.join(series_names)}" if len(series_names) > 1 else ""
                
                summary = f"""Key Insights - {title}

• Data spans from {years[0]} to {years[-1]} ({len(years)} data points)
• Starting value: {start_val:,.0f} {unit} ({years[0]})
• Ending value: {end_val:,.0f} {unit} ({years[-1]})
• Overall change: {growth:+.0f}% over the period
{series_summary}
• Trend shows {'growth' if growth > 0 else 'decline' if growth < 0 else 'stability'} pattern"""
            else:
                summary = f"Key Insights - {title}\n\n• Multi-scenario analysis\n• Unit: {unit}\n• Data from {years[0]} to {years[-1]}\n• {len(series_names)} scenarios compared"
        else:
            summary = f"Key Insights - {title}\n\n• Chart shows trend data\n• Unit: {unit}\n• [Add specific insights here]"
    
    elif chart_info['type'] == 'stacked':
        years = chart_info['years']
        totals = chart_info['totals']
        series_names = chart_info['series_names']

        if years and len(years) >= 2:
            latest_year = max(years)
            latest_total = totals[latest_year]
            earliest_total = totals[min(years)]
            growth = ((latest_total - earliest_total) / earliest_total * 100) if earliest_total > 0 else 0

            series_info = []
            if latest_year in chart_info['series']:
                for series in series_names:
                    value = chart_info['series'][latest_year].get(series, 0)
                    percentage = (value / latest_total * 100) if latest_total > 0 else 0
                    series_info.append(f"{series}: {percentage:.1f}%")

            summary = f"""Key Insights - {title}

• Total capacity: {latest_total:,.0f} {unit} in {latest_year}
• Market composition: {' | '.join(series_info) if series_info else '[Add breakdown here]'}
• Period growth: {growth:+.0f}% from {min(years)} to {max(years)}
• Data covers {len(years)} years of market development
• [Add additional insights here]"""
        else:
            summary = f"Key Insights - {title}\n\n• Stacked chart analysis\n• Unit: {unit}\n• [Add specific insights here]"

    elif chart_info['type'] == 'pie':
        categories = chart_info['categories']
        values = chart_info['values']
        total = chart_info['total']

        if categories and values:
            # Calculate percentages for each category
            percentages = [(value / total * 100) if total > 0 else 0 for value in values]

            # Find largest and smallest segments
            if len(values) > 0:
                max_idx = values.index(max(values))
                largest_category = categories[max_idx]
                largest_percentage = percentages[max_idx]

                # Create breakdown list
                breakdown = []
                for i, (cat, pct) in enumerate(zip(categories, percentages)):
                    if i < 5:  # Show top 5 categories
                        breakdown.append(f"{cat}: {pct:.1f}%")

                summary = f"""Key Insights - {title}

• Total: {total:,.0f} {unit}
• Number of categories: {len(categories)}
• Largest segment: {largest_category} ({largest_percentage:.1f}%)
• Distribution: {' | '.join(breakdown)}
• [Add additional context here]"""
            else:
                summary = f"Key Insights - {title}\n\n• Pie chart showing distribution\n• Unit: {unit}\n• [Add specific insights here]"
        else:
            summary = f"Key Insights - {title}\n\n• Pie chart analysis\n• Unit: {unit}\n• [Add specific insights here]"

    else:
        summary = f"Key Insights - {title}\n\n• Data visualization\n• Unit: {unit}\n• [Customize this content]\n• [Add relevant insights]\n• [Modify as needed]"
    
    return summary

def create_chart_from_data(slide, plot_data, chart_info):
    """Create appropriate chart based on data type"""
    
    plot_shape_position = None
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and '[plot]' in shape.text:
            plot_shape_position = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            shapes_to_remove.append(shape)
            break
    
    for shape in shapes_to_remove:
        slide.shapes._spTree.remove(shape._element)
    
    if not plot_shape_position:
        plot_shape_position = {
            'left': Inches(1),
            'top': Inches(2),
            'width': Inches(7),
            'height': Inches(4)
        }
    
    if chart_info['type'] == 'line':
        return create_line_chart(slide, plot_shape_position, chart_info)
    elif chart_info['type'] == 'stacked':
        return create_stacked_chart(slide, plot_shape_position, chart_info)
    elif chart_info['type'] == 'pie':
        return create_pie_chart(slide, plot_shape_position, chart_info)
    else:
        return create_basic_chart(slide, plot_shape_position, chart_info)

def create_line_chart(slide, position, chart_info):
    """Create modern, minimal line chart with multiple series support"""
    chart_data_pptx = CategoryChartData()
    chart_data_pptx.categories = chart_info['categories']
    
    # Add all series to the chart
    for series_name, values in chart_info['series'].items():
        chart_data_pptx.add_series(series_name, values)
    
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        position['left'], position['top'],
        position['width'], position['height'],
        chart_data_pptx
    )
    
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True if len(chart_info['series']) > 1 else False
    
    # Position legend at bottom for multi-series
    if chart.has_legend:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = RGBColor(80, 80, 80)  # Dark gray
    
    # Modern minimal styling
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(240, 240, 240)  # Very light gray
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    
    # Remove minor gridlines
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    
    # Style each line with different colors for multi-series
    line_colors = [
        RGBColor(64, 150, 255),   # Blue for first series
        RGBColor(255, 140, 60),   # Orange for second series
        RGBColor(80, 200, 120),   # Green for third series
        RGBColor(200, 80, 200),   # Purple for fourth series
        RGBColor(255, 80, 80),    # Red for fifth series
    ]
    
    for i, series in enumerate(chart.series):
        color = line_colors[i % len(line_colors)]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(3)  # Thicker line
        
        # Add subtle markers
        series.marker.style = 2  # Circle
        series.marker.size = 6
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.marker.format.line.color.rgb = RGBColor(255, 255, 255)
        series.marker.format.line.width = Pt(1)
    
    # Modern, light axis styling
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(120, 120, 120)  # Light gray
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(120, 120, 120)
    
    # Remove axis lines for cleaner look
    chart.category_axis.format.line.color.rgb = RGBColor(230, 230, 230)
    chart.value_axis.format.line.color.rgb = RGBColor(230, 230, 230)
    
    return chart_shape

def create_stacked_chart(slide, position, chart_info):
    """Create modern, minimal stacked column chart"""
    chart_data_pptx = CategoryChartData()
    chart_data_pptx.categories = chart_info['categories']
    
    series_names = chart_info['series_names']
    for series in series_names:
        values = []
        for year in chart_info['years']:
            values.append(chart_info['series'][year].get(series, 0))
        chart_data_pptx.add_series(series, values)
    
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        position['left'], position['top'],
        position['width'], position['height'],
        chart_data_pptx
    )
    
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = RGBColor(80, 80, 80)  # Dark gray
    
    # Modern color palette - subtle, professional
    modern_colors = [
        RGBColor(88, 166, 255),   # Bright blue for primary series
        RGBColor(155, 200, 255),  # Lighter blue for secondary series
        RGBColor(220, 235, 255),  # Very light blue for additional series
    ]
    
    # Remove borders completely by setting gap width and overlap
    try:
        chart.plots[0].gap_width = 150  # Reduce gap between categories
        chart.plots[0].overlap = 100    # Make bars overlap to remove white lines
    except:
        pass
    
    for i, series in enumerate(chart.series):
        if i < len(modern_colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = modern_colors[i]
            # Multiple approaches to remove white lines
            try:
                # Method 1: Set line to background
                series.format.line.fill.background()
            except:
                pass
            try:
                # Method 2: Set line color to match fill and make very thin
                series.format.line.color.rgb = modern_colors[i]
                series.format.line.width = Pt(0.1)
            except:
                pass
    
    # Minimal gridlines
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(240, 240, 240)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    chart.value_axis.has_minor_gridlines = False
    
    # Clean axes
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    
    # Light, minimal axis text
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(120, 120, 120)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(120, 120, 120)
    
    # Subtle axis lines
    chart.category_axis.format.line.color.rgb = RGBColor(230, 230, 230)
    chart.value_axis.format.line.color.rgb = RGBColor(230, 230, 230)
    
    return chart_shape

def create_pie_chart(slide, position, chart_info):
    """Create modern, minimal pie chart"""
    chart_data_pptx = CategoryChartData()
    chart_data_pptx.categories = chart_info['categories']

    # Add single series for pie chart
    chart_data_pptx.add_series('', chart_info['values'])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        position['left'], position['top'],
        position['width'], position['height'],
        chart_data_pptx
    )

    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = RGBColor(80, 80, 80)  # Dark gray

    # Modern color palette for pie slices
    modern_colors = [
        RGBColor(88, 166, 255),   # Bright blue
        RGBColor(255, 140, 60),   # Orange
        RGBColor(80, 200, 120),   # Green
        RGBColor(200, 80, 200),   # Purple
        RGBColor(255, 80, 80),    # Red
        RGBColor(255, 200, 60),   # Yellow
        RGBColor(100, 180, 180),  # Teal
        RGBColor(180, 100, 180),  # Magenta
    ]

    # Style each pie slice
    for i, point in enumerate(chart.series[0].points):
        color = modern_colors[i % len(modern_colors)]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

        # Remove borders for clean look
        try:
            point.format.line.fill.background()
        except:
            pass

    # Show data labels with percentages
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = RGBColor(60, 60, 60)
    data_labels.number_format = '0.0%'  # Show as percentage
    data_labels.position = 2  # XL_LABEL_POSITION.OUTSIDE_END

    return chart_shape

def create_basic_chart(slide, position, chart_info):
    """Fallback basic chart creation"""
    text_box = slide.shapes.add_textbox(
        position['left'], position['top'],
        position['width'], position['height']
    )
    text_frame = text_box.text_frame
    text_frame.text = "Chart data found but type not fully supported.\nPlease customize this area with your chart."
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return text_box

def update_slide_with_plot_data(slide, plot_data, slide_number):
    """Update an existing slide with plot data using modern, minimal styling"""
    placeholders_updated = {'title': False, 'plot': False, 'summary': False}
    
    chart_info = prepare_chart_data_generic(plot_data)
    summary_content = generate_generic_content(plot_data, chart_info)
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            shape_text = shape.text.strip()
            
            if ('[plot title]' in shape_text or 
                'Click to add title' in shape_text or 
                shape_text == '[plot title]'):
                
                shape.text = plot_data.get('title', 'Data Analysis Chart')
                
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(22)  # Slightly smaller
                    p.font.bold = False  # Not bold for modern look
                    p.font.name = "Segoe UI Light"  # Modern, light font
                    p.font.color.rgb = RGBColor(60, 60, 60)  # Dark gray instead of black
                placeholders_updated['title'] = True
            
            elif ('[plot summary]' in shape_text or 
                  'Click to add subtitle' in shape_text or
                  shape_text == '[plot summary]'):
                
                shape.text = summary_content
                
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.word_wrap = True
                    for i, paragraph in enumerate(shape.text_frame.paragraphs):
                        if i == 0:  # Header
                            paragraph.font.size = Pt(13)  # Smaller header
                            paragraph.font.bold = False  # Not bold
                            paragraph.font.name = "Segoe UI"  # Modern font
                            paragraph.font.color.rgb = RGBColor(80, 80, 80)  # Medium gray
                        else:  # Bullet points
                            paragraph.font.size = Pt(11)  # Smaller text
                            paragraph.font.bold = False
                            paragraph.font.name = "Segoe UI Light"  # Light font
                            paragraph.font.color.rgb = RGBColor(100, 100, 100)  # Light gray
                        paragraph.space_after = Pt(8)  # More spacing
                placeholders_updated['summary'] = True
            
            elif 'company name' in shape_text.lower():
                if slide_number > 1:
                    current_text = shape.text
                    # Add slide number to company name if not already present
                    updated_text = re.sub(r'\s+\d+\s*$', f"    {slide_number}", current_text)
                    if updated_text == current_text:
                        updated_text = f"{current_text.rstrip()}    {slide_number}"
                    shape.text = updated_text
                
                # Style footer text
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    p.font.size = Pt(9)
                    p.font.bold = False
                    p.font.name = "Segoe UI Light"
                    p.font.color.rgb = RGBColor(140, 140, 140)  # Very light gray
    
    chart_shape = create_chart_from_data(slide, plot_data, chart_info)
    if chart_shape:
        placeholders_updated['plot'] = True
    
    return placeholders_updated

def create_powerpoint_from_json_all_plots(template_path, json_file_path, output_path):
    """Create PowerPoint presentation using multi-slide template"""
    
    plot_items, full_data = load_json_data(json_file_path)
    
    print(f"Found {len(plot_items)} plots in JSON file")
    for i, plot_item in enumerate(plot_items):
        title = plot_item['payload'].get('title', f'Plot {i+1}')
        print(f"  {i+1}. {title}")
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    # Load template with multiple slides
    prs = Presentation(template_path)
    
    print(f"Template has {len(prs.slides)} slides")
    
    # Check if we have enough slides in template
    if len(prs.slides) < len(plot_items):
        print(f"Warning: Template has {len(prs.slides)} slides but need {len(plot_items)}")
        print("Will only update available slides")
    
    # Update each slide with corresponding plot data
    slides_updated = 0
    for i, plot_item in enumerate(plot_items):
        if i >= len(prs.slides):
            print(f"Skipping plot {i+1} - no more slides in template")
            break
            
        plot_data = plot_item['payload']
        title = plot_data.get('title', f'Plot {i+1}')
        slide = prs.slides[i]
        
        print(f"Updating slide {i+1}: {title}")
        
        # Update this slide with plot data
        updates = update_slide_with_plot_data(slide, plot_data, i+1)
        slides_updated += 1
        
        status_symbols = {"success": 0, "failed": 0}
        for placeholder, updated in updates.items():
            if updated:
                status_symbols["success"] += 1
            else:
                status_symbols["failed"] += 1
        
        print(f"   Updated: {status_symbols['success']} placeholders, {status_symbols['failed']} not found")
    
    # Remove unused slides if template has more slides than plots
    if len(prs.slides) > len(plot_items):
        slides_to_remove = []
        for i in range(len(plot_items), len(prs.slides)):
            slides_to_remove.append(i)
        
        # Remove slides in reverse order to avoid index issues
        for slide_index in reversed(slides_to_remove):
            rId = prs.slides._sldIdLst[slide_index].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slide_index]
        
        print(f"Removed {len(slides_to_remove)} unused slides")
    
    # Save presentation
    prs.save(output_path)
    print(f"SUCCESS! Generated: {output_path}")
    print(f"Updated {slides_updated} slides using your template")
    
    return output_path

# Main execution for multi-slide template
if __name__ == "__main__":
    CONFIG = {
        'json_file': "conversation_110_2025-09-22_11.json",  # Updated to use the multi-series data
        'template_file': "template.pptx",  # Your multi-slide template
        'output_file': "final_presentation.pptx"
    }
    
    print("MULTI-SLIDE TEMPLATE POWERPOINT GENERATOR")
    print("="*60)
    print("Uses your multi-slide template - no duplication needed!")
    
    try:
        result = create_powerpoint_from_json_all_plots(
            CONFIG['template_file'],
            CONFIG['json_file'],
            CONFIG['output_file']
        )
        
        if result:
            print("\nSUCCESS! Your presentation is ready.")
            print(f"File: {CONFIG['output_file']}")
            print("\nThis approach:")
            print("• Uses your existing multi-slide template")
            print("• Updates each slide with plot data")
            print("• Preserves all template design elements")
            print("• Supports multi-series line charts")
            print("• No slide duplication issues!")
        
    except Exception as e:
        print(f"Error: {e}")
    
    print("="*60)